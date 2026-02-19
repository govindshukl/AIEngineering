#!/usr/bin/env python3
"""
AI/ML Center of Excellence Charter - PowerPoint Generator
Creates a professional presentation from the CoE Charter document.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path


class CoEPresentationGenerator:
    """Generates a PowerPoint presentation for the CoE Charter."""

    def __init__(self, output_path):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        self.prs.slide_height = Inches(7.5)
        self.output_path = Path(output_path)

        # Brand colors
        self.PRIMARY_BLUE = RGBColor(0, 51, 102)
        self.SECONDARY_BLUE = RGBColor(0, 102, 153)
        self.ACCENT_ORANGE = RGBColor(237, 125, 49)
        self.LIGHT_GRAY = RGBColor(240, 240, 240)
        self.DARK_GRAY = RGBColor(64, 64, 64)
        self.WHITE = RGBColor(255, 255, 255)

    def add_title_slide(self, title, subtitle):
        """Add a title slide."""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # Background shape
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.PRIMARY_BLUE
        shape.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.WHITE
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = self.ACCENT_ORANGE
        p.alignment = PP_ALIGN.CENTER

        # Footer
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Bank ABC | Innovation & Digitization Department"
        p.font.size = Pt(14)
        p.font.color.rgb = self.WHITE
        p.alignment = PP_ALIGN.CENTER

    def add_section_slide(self, section_number, section_title):
        """Add a section divider slide."""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # Left accent bar
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.5), self.prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.PRIMARY_BLUE
        shape.line.fill.background()

        # Section number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1), Inches(2.5), Inches(1.5), Inches(1.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = self.ACCENT_ORANGE
        circle.line.fill.background()

        # Number text
        num_box = slide.shapes.add_textbox(Inches(1), Inches(2.7), Inches(1.5), Inches(1))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(section_number)
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.WHITE
        p.alignment = PP_ALIGN.CENTER

        # Section title
        title_box = slide.shapes.add_textbox(Inches(3), Inches(2.7), Inches(9), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = section_title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.PRIMARY_BLUE

    def add_content_slide(self, title, bullets, subtitle=None):
        """Add a content slide with bullet points."""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # Header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.PRIMARY_BLUE
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.WHITE

        # Subtitle if provided
        start_y = Inches(1.5)
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(0.5))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(16)
            p.font.italic = True
            p.font.color.rgb = self.DARK_GRAY
            start_y = Inches(1.9)

        # Bullet points
        bullet_box = slide.shapes.add_textbox(Inches(0.5), start_y, Inches(12), Inches(5))
        tf = bullet_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(20)
            p.font.color.rgb = self.DARK_GRAY
            p.space_after = Pt(12)

    def add_two_column_slide(self, title, left_content, right_content, left_title="", right_title=""):
        """Add a slide with two columns."""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # Header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.PRIMARY_BLUE
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.WHITE

        # Left column title
        if left_title:
            left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.5), Inches(0.5))
            tf = left_title_box.text_frame
            p = tf.paragraphs[0]
            p.text = left_title
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.SECONDARY_BLUE

        # Left column content
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(5.5), Inches(5))
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(left_content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(16)
            p.font.color.rgb = self.DARK_GRAY
            p.space_after = Pt(8)

        # Right column title
        if right_title:
            right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(5.5), Inches(0.5))
            tf = right_title_box.text_frame
            p = tf.paragraphs[0]
            p.text = right_title
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.SECONDARY_BLUE

        # Right column content
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.9), Inches(5.5), Inches(5))
        tf = right_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(right_content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(16)
            p.font.color.rgb = self.DARK_GRAY
            p.space_after = Pt(8)

    def add_table_slide(self, title, headers, rows):
        """Add a slide with a table."""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # Header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.PRIMARY_BLUE
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.WHITE

        # Table
        num_rows = len(rows) + 1
        num_cols = len(headers)
        table_width = Inches(12)
        table_height = Inches(min(5.5, 0.5 + len(rows) * 0.45))

        table = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(0.5), Inches(1.5),
            table_width, table_height
        ).table

        # Set column widths
        col_width = table_width / num_cols
        for i in range(num_cols):
            table.columns[i].width = int(col_width)

        # Header row
        for i, header_text in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.PRIMARY_BLUE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = self.WHITE
                paragraph.alignment = PP_ALIGN.CENTER

        # Data rows
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_text in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_text)
                # Alternate row colors
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.LIGHT_GRAY
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                    paragraph.font.color.rgb = self.DARK_GRAY

    def add_org_chart_slide(self, title):
        """Add an organizational structure slide."""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # Header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.PRIMARY_BLUE
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.WHITE

        # CDO Box
        cdo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.5), Inches(2.9), Inches(0.6))
        cdo.fill.solid()
        cdo.fill.fore_color.rgb = self.DARK_GRAY
        cdo.text_frame.paragraphs[0].text = "Chief Digital Officer"
        cdo.text_frame.paragraphs[0].font.size = Pt(11)
        cdo.text_frame.paragraphs[0].font.color.rgb = self.WHITE
        cdo.text_frame.paragraphs[0].font.bold = True
        cdo.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Steering Committee Box
        steer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.7), Inches(2.3), Inches(3.9), Inches(0.6))
        steer.fill.solid()
        steer.fill.fore_color.rgb = self.SECONDARY_BLUE
        steer.text_frame.paragraphs[0].text = "AI/ML Steering Committee"
        steer.text_frame.paragraphs[0].font.size = Pt(11)
        steer.text_frame.paragraphs[0].font.color.rgb = self.WHITE
        steer.text_frame.paragraphs[0].font.bold = True
        steer.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # CoE Box
        coe = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(3.1), Inches(4.9), Inches(0.7))
        coe.fill.solid()
        coe.fill.fore_color.rgb = self.PRIMARY_BLUE
        coe.text_frame.paragraphs[0].text = "AI/ML Center of Excellence (5 FTEs)"
        coe.text_frame.paragraphs[0].font.size = Pt(12)
        coe.text_frame.paragraphs[0].font.color.rgb = self.WHITE
        coe.text_frame.paragraphs[0].font.bold = True
        coe.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Core team boxes
        roles = ["CoE Lead", "AI Architect", "Model Risk\nSpecialist", "AI Product\nManager", "AI Trainer"]
        start_x = 1.0
        for i, role in enumerate(roles):
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(start_x + i * 2.3), Inches(4.0),
                Inches(2.1), Inches(0.7)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.ACCENT_ORANGE
            tf = box.text_frame
            tf.paragraphs[0].text = role
            tf.paragraphs[0].font.size = Pt(10)
            tf.paragraphs[0].font.color.rgb = self.WHITE
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.word_wrap = True

        # Champion boxes
        champions = ["Retail\nChampion", "Corporate\nChampion", "Treasury\nChampion", "Ops\nChampion", "Risk\nChampion", "Digital\nChampion"]
        start_x = 0.5
        for i, champ in enumerate(champions):
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(start_x + i * 2.05), Inches(5.2),
                Inches(1.85), Inches(0.65)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(144, 164, 174)  # Light blue-gray
            tf = box.text_frame
            tf.paragraphs[0].text = champ
            tf.paragraphs[0].font.size = Pt(9)
            tf.paragraphs[0].font.color.rgb = self.WHITE
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.word_wrap = True

        # Labels
        label1 = slide.shapes.add_textbox(Inches(0.5), Inches(3.85), Inches(2), Inches(0.3))
        label1.text_frame.paragraphs[0].text = "Core Team:"
        label1.text_frame.paragraphs[0].font.size = Pt(10)
        label1.text_frame.paragraphs[0].font.bold = True
        label1.text_frame.paragraphs[0].font.color.rgb = self.DARK_GRAY

        label2 = slide.shapes.add_textbox(Inches(0.5), Inches(5.05), Inches(3), Inches(0.3))
        label2.text_frame.paragraphs[0].text = "Federated Champions (dotted-line):"
        label2.text_frame.paragraphs[0].font.size = Pt(10)
        label2.text_frame.paragraphs[0].font.bold = True
        label2.text_frame.paragraphs[0].font.color.rgb = self.DARK_GRAY

        # BU Teams
        bu_label = slide.shapes.add_textbox(Inches(4.5), Inches(6.1), Inches(4.5), Inches(0.4))
        bu_label.text_frame.paragraphs[0].text = "Business Unit AI Teams"
        bu_label.text_frame.paragraphs[0].font.size = Pt(12)
        bu_label.text_frame.paragraphs[0].font.italic = True
        bu_label.text_frame.paragraphs[0].font.color.rgb = self.DARK_GRAY
        bu_label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def add_kpi_slide(self, title, kpis):
        """Add a KPI dashboard slide."""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # Header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.PRIMARY_BLUE
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.WHITE

        # KPI cards
        colors = [self.PRIMARY_BLUE, self.SECONDARY_BLUE, self.ACCENT_ORANGE, RGBColor(76, 175, 80)]
        for i, (kpi_title, kpi_target, kpi_desc) in enumerate(kpis[:4]):
            row = i // 2
            col = i % 2
            x = Inches(0.75 + col * 6.2)
            y = Inches(1.6 + row * 2.8)

            # Card background
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.8), Inches(2.4)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self.LIGHT_GRAY
            card.line.color.rgb = colors[i]
            card.line.width = Pt(3)

            # KPI Title
            kpi_title_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(5.4), Inches(0.5))
            tf = kpi_title_box.text_frame
            p = tf.paragraphs[0]
            p.text = kpi_title
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = colors[i]

            # KPI Target
            kpi_target_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.7), Inches(5.4), Inches(0.8))
            tf = kpi_target_box.text_frame
            p = tf.paragraphs[0]
            p.text = kpi_target
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = self.DARK_GRAY

            # KPI Description
            kpi_desc_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.5), Inches(5.4), Inches(0.7))
            tf = kpi_desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = kpi_desc
            p.font.size = Pt(12)
            p.font.color.rgb = self.DARK_GRAY

    def generate(self):
        """Generate the full presentation."""
        print("Generating CoE Charter PowerPoint presentation...")

        # Slide 1: Title
        self.add_title_slide(
            "AI/ML Center of Excellence",
            "Charter & Operating Model"
        )

        # Slide 2: Agenda
        self.add_content_slide(
            "Agenda",
            [
                "Executive Summary & Strategic Mandate",
                "Mission, Vision & Objectives",
                "Scope & Ownership Areas",
                "Organizational Structure",
                "Core Team Roles & Responsibilities",
                "Federated Champion Network",
                "Services Catalog & SLAs",
                "Success Metrics & KPIs",
                "Evolution Roadmap"
            ]
        )

        # Slide 3: Section - Executive Summary
        self.add_section_slide(1, "Executive Summary")

        # Slide 4: Charter Statement
        self.add_content_slide(
            "Charter Statement",
            [
                "Establishes the AI/ML CoE as the central governing body for all AI, ML, and Intelligent Automation initiatives",
                "Operates under Innovation & Digitization Department",
                "Mandate to accelerate responsible AI adoption while ensuring regulatory compliance",
                "Empowered to set standards, enable innovation, manage risk, build capability, and drive value"
            ],
            subtitle="Formalizing AI governance and operational excellence"
        )

        # Slide 5: Operating Principles
        self.add_two_column_slide(
            "Operating Principles",
            [
                "Lean Core, Federated Reach: Small central team with embedded champions",
                "Enabler, Not Bottleneck: Streamlined processes with tiered approval",
                "Responsible AI: Ethics, fairness, transparency embedded"
            ],
            [
                "Business-First: Technology decisions driven by business value",
                "Continuous Learning: Adaptive approach with regular retrospectives",
                "Risk-Based Governance: Proportionate controls by tier"
            ],
            "Core Principles",
            "Operational Philosophy"
        )

        # Slide 6: Section - Mission & Vision
        self.add_section_slide(2, "Mission & Vision")

        # Slide 7: Mission Statement
        self.add_content_slide(
            "Mission Statement",
            [
                "Enable Bank ABC to become an AI-driven organization",
                "Establish world-class governance frameworks",
                "Provide scalable infrastructure and reusable assets",
                "Build organization-wide AI capability and literacy",
                "Ensure full regulatory compliance and responsible AI practices"
            ],
            subtitle="Our Purpose"
        )

        # Slide 8: Strategic Objectives
        self.add_table_slide(
            "Strategic Objectives",
            ["Objective", "Target", "Timeline"],
            [
                ["Establish governance framework", "12 core documents approved", "Q1 2025"],
                ["Launch AI/Automation Hub", "Platform operational", "Q2 2025"],
                ["Certify AI practitioners", "100% of practitioners", "Q4 2025"],
                ["Production AI use cases", "30+ active models", "Q4 2026"],
                ["Reduce deployment time", "50% reduction", "Q4 2026"]
            ]
        )

        # Slide 9: Section - Scope
        self.add_section_slide(3, "Scope & Ownership")

        # Slide 10: Ownership Areas
        self.add_two_column_slide(
            "CoE Ownership Areas",
            [
                "AI Governance & Policy: Framework, ethics, compliance",
                "Model Risk Management: Inventory, validation, monitoring",
                "Reusable Assets & MLOps: Hub, pipelines, templates"
            ],
            [
                "Copilot & GenAI Governance: Policies, prompts, controls",
                "Specialized AI Domains: Fraud, Agentic AI, high-risk",
                "AI Capability Building: Training, certification, CoP"
            ],
            "Primary Ownership",
            "Specialized Areas"
        )

        # Slide 11: Section - Organization
        self.add_section_slide(4, "Organizational Structure")

        # Slide 12: Org Chart
        self.add_org_chart_slide("CoE Organizational Structure")

        # Slide 13: Core Team Composition
        self.add_table_slide(
            "Core Team Composition (5 FTEs)",
            ["Role", "FTE", "Primary Focus"],
            [
                ["CoE Lead", "1.0", "Strategy, governance, stakeholder management"],
                ["AI Architect", "1.0", "Technical standards, MLOps, platform"],
                ["Model Risk Specialist", "1.0", "Validation, inventory, compliance"],
                ["AI Product Manager", "1.0", "Hub management, use case intake"],
                ["AI Trainer/Enabler", "0.5-1.0", "Literacy programs, champion network"]
            ]
        )

        # Slide 14: Section - Champions
        self.add_section_slide(5, "Federated Champion Network")

        # Slide 15: Champion Model
        self.add_content_slide(
            "Federated Champion Model",
            [
                "BU-embedded AI advocates with dotted-line reporting to CoE",
                "Minimum 20% allocation to Champion responsibilities",
                "Selection criteria: Senior level, domain expertise, technical aptitude",
                "Local approval authority for Tier 3 initiatives",
                "First-line support for Copilot and Hub users",
                "Bi-weekly sync with CoE, quarterly workshops"
            ],
            subtitle="Extending CoE reach without bottlenecks"
        )

        # Slide 16: Champion Responsibilities
        self.add_two_column_slide(
            "Champion Responsibilities",
            [
                "Collect and prioritize AI use cases from BU",
                "Submit use cases through intake process",
                "Ensure BU compliance with CoE standards",
                "Conduct local design reviews (Tier 3)"
            ],
            [
                "Cascade CoE training to BU teams",
                "First-line support for Copilot users",
                "Report challenges and improvements to CoE",
                "Participate in CoP and knowledge sharing"
            ],
            "Demand & Governance",
            "Enablement & Support"
        )

        # Slide 17: Section - Services
        self.add_section_slide(6, "Services Catalog")

        # Slide 18: Services Overview
        self.add_table_slide(
            "Services Catalog Overview",
            ["Category", "Key Services", "Delivery"],
            [
                ["Governance", "Use case intake, risk classification, gate reviews", "CoE Lead / Product Manager"],
                ["Technical", "Architecture review, MLOps onboarding, templates", "AI Architect"],
                ["Risk Management", "Model inventory, validation coordination, monitoring", "Model Risk Specialist"],
                ["Enablement", "Training programs, certification, CoP facilitation", "AI Trainer"],
                ["Platform", "Hub access, Copilot provisioning, usage analytics", "AI Product Manager"]
            ]
        )

        # Slide 19: SLA Commitments
        self.add_table_slide(
            "SLA Commitments",
            ["Service", "SLA", "Measurement"],
            [
                ["Use case intake response", "2 business days", "Time to initial triage"],
                ["Tier classification", "5 business days", "Time to risk assessment"],
                ["Architecture review (Tier 3)", "5 business days", "Time to approval"],
                ["Architecture review (Tier 2)", "10 business days", "Time to approval"],
                ["Champion support response", "1 business day", "Time to acknowledge"],
                ["Production incident (P1)", "4 hours", "Time to engage"]
            ]
        )

        # Slide 20: Section - KPIs
        self.add_section_slide(7, "Success Metrics")

        # Slide 21: Strategic KPIs
        self.add_kpi_slide(
            "Strategic KPIs",
            [
                ("Active AI Use Cases", "30+", "Production models by Q4 2026"),
                ("AI Practitioner Certification", "100%", "All practitioners certified by Q4 2025"),
                ("Regulatory Findings", "0 Critical", "Zero critical AI-related findings"),
                ("Process Automation", "40%", "Coverage target by Year 2")
            ]
        )

        # Slide 22: Operational KPIs
        self.add_table_slide(
            "Operational KPIs",
            ["KPI", "Target", "Owner"],
            [
                ["Intake-to-production (Tier 3)", "< 8 weeks", "AI Product Manager"],
                ["Intake-to-production (Tier 2)", "< 16 weeks", "AI Product Manager"],
                ["Model validation pass rate", "> 80%", "Model Risk Specialist"],
                ["MLOps pipeline availability", "> 99.5%", "AI Architect"],
                ["Hub monthly active users", "200+", "AI Product Manager"],
                ["Training satisfaction score", "> 4.2 / 5.0", "AI Trainer"]
            ]
        )

        # Slide 23: Section - Roadmap
        self.add_section_slide(8, "Evolution Roadmap")

        # Slide 24: Maturity Model
        self.add_table_slide(
            "CoE Maturity Model",
            ["Stage", "Characteristics", "Timeline"],
            [
                ["1. Foundation", "Governance framework, core team, basic platform", "Q1 2025"],
                ["2. Enablement", "Hub launched, Champions active, training programs", "Q2-Q3 2025"],
                ["3. Scale", "20+ production use cases, self-service dominant", "Q4 2025 - Q2 2026"],
                ["4. Optimization", "MLOps Level 3, automated governance", "Q3-Q4 2026"],
                ["5. Innovation", "Leading edge AI, industry leadership", "2027+"]
            ]
        )

        # Slide 25: Capability Roadmap
        self.add_table_slide(
            "Capability Roadmap 2025",
            ["Quarter", "Governance", "Platform", "Capability"],
            [
                ["Q1 2025", "Framework v1.0 approved", "MLOps foundation", "Core team hired"],
                ["Q2 2025", "Gate process operational", "Hub v1.0 launch", "Champion network live"],
                ["Q3 2025", "Model inventory populated", "Feature store v1.0", "Level 1-2 training"],
                ["Q4 2025", "Automated compliance", "Hub v2.0 self-service", "Certification program"]
            ]
        )

        # Slide 26: Summary
        self.add_content_slide(
            "Summary: Key Takeaways",
            [
                "CoE established as central governing body for AI/ML and Automation",
                "Lean core team (5 FTEs) with federated Champion network",
                "Clear ownership of governance, platform, risk, and enablement",
                "Tiered service model: Strategic, Standard, Self-Service",
                "Defined SLAs and success metrics for accountability",
                "Phased evolution from Foundation to Innovation leadership"
            ],
            subtitle="Enabling responsible AI adoption at scale"
        )

        # Slide 27: Next Steps
        self.add_content_slide(
            "Next Steps",
            [
                "Obtain charter approval from CDO and key stakeholders",
                "Finalize core team hiring and onboarding",
                "Identify and appoint Business Unit Champions",
                "Launch AI/Automation Hub (Q2 2025)",
                "Roll out AI Awareness training (Level 1)",
                "Begin use case intake and pipeline building"
            ],
            subtitle="Immediate actions for CoE establishment"
        )

        # Slide 28: Q&A
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # Background
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.PRIMARY_BLUE
        shape.line.fill.background()

        # Q&A text
        qa_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5))
        tf = qa_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Questions & Discussion"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.WHITE
        p.alignment = PP_ALIGN.CENTER

        # Contact info
        contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(1))
        tf = contact_box.text_frame
        p = tf.paragraphs[0]
        p.text = "AI Engineering Team | Innovation & Digitization"
        p.font.size = Pt(20)
        p.font.color.rgb = self.ACCENT_ORANGE
        p.alignment = PP_ALIGN.CENTER

        # Save
        self.prs.save(self.output_path)
        print(f"✓ Presentation saved: {self.output_path}")
        print(f"  File size: {self.output_path.stat().st_size / 1024:.1f} KB")
        print(f"  Total slides: {len(self.prs.slides)}")


def main():
    """Main function to generate the presentation."""
    output_path = Path('/Users/govind/AIEngineering/governance-docs-word/12-AI-ML-Center-of-Excellence-Charter.pptx')

    generator = CoEPresentationGenerator(output_path)
    generator.generate()

    print("\nPresentation generation complete!")


if __name__ == '__main__':
    main()
