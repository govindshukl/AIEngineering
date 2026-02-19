#!/usr/bin/env python3
"""
AI/ML Center of Excellence - Executive Summary PowerPoint
Concise presentation for executive management (10-12 slides).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path


class ExecPresentationGenerator:
    """Generates an executive summary presentation."""

    def __init__(self, output_path):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.output_path = Path(output_path)

        # Brand colors
        self.PRIMARY = RGBColor(0, 51, 102)
        self.SECONDARY = RGBColor(0, 102, 153)
        self.ACCENT = RGBColor(237, 125, 49)
        self.SUCCESS = RGBColor(76, 175, 80)
        self.LIGHT = RGBColor(245, 245, 245)
        self.DARK = RGBColor(51, 51, 51)
        self.WHITE = RGBColor(255, 255, 255)

    def add_title_slide(self):
        """Slide 1: Title."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Full background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.PRIMARY
        bg.line.fill.background()

        # Title
        title = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.2))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "AI/ML Center of Excellence"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.WHITE
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.8))
        tf = sub.text_frame
        p = tf.paragraphs[0]
        p.text = "Executive Summary"
        p.font.size = Pt(28)
        p.font.color.rgb = self.ACCENT
        p.alignment = PP_ALIGN.CENTER

        # Footer
        footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5))
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.text = "Bank ABC | Innovation & Digitization | January 2025"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(180, 180, 180)
        p.alignment = PP_ALIGN.CENTER

    def add_header(self, slide, title):
        """Add standard header to slide."""
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.1))
        header.fill.solid()
        header.fill.fore_color.rgb = self.PRIMARY
        header.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = self.WHITE

    def add_why_slide(self):
        """Slide 2: Why a CoE?"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_header(slide, "Why an AI/ML Center of Excellence?")

        # Left side - Challenges
        challenge_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.5))
        tf = challenge_title.text_frame
        p = tf.paragraphs[0]
        p.text = "Current Challenges"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(192, 57, 43)

        challenges = [
            "Fragmented AI initiatives across BUs",
            "Inconsistent governance & risk management",
            "Regulatory scrutiny increasing (CBB)",
            "Duplicated efforts, no reusable assets",
            "Skills gap limiting AI adoption"
        ]

        y_pos = 1.9
        for challenge in challenges:
            box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(5.8), Inches(0.45))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = "✗  " + challenge
            p.font.size = Pt(16)
            p.font.color.rgb = self.DARK
            y_pos += 0.55

        # Right side - Solution
        solution_title = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(5.8), Inches(0.5))
        tf = solution_title.text_frame
        p = tf.paragraphs[0]
        p.text = "CoE Solution"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.SUCCESS

        solutions = [
            "Central governance with clear accountability",
            "Standardized risk framework (Tier 1/2/3)",
            "Regulatory-aligned processes (CBB, Basel)",
            "Shared platform, reusable components",
            "Structured training & capability building"
        ]

        y_pos = 1.9
        for solution in solutions:
            box = slide.shapes.add_textbox(Inches(6.8), Inches(y_pos), Inches(5.8), Inches(0.45))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = "✓  " + solution
            p.font.size = Pt(16)
            p.font.color.rgb = self.DARK
            y_pos += 0.55

        # Bottom callout
        callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.2))
        callout.fill.solid()
        callout.fill.fore_color.rgb = self.LIGHT
        callout.line.color.rgb = self.ACCENT
        callout.line.width = Pt(2)

        callout_text = slide.shapes.add_textbox(Inches(0.7), Inches(5.7), Inches(12), Inches(0.8))
        tf = callout_text.text_frame
        p = tf.paragraphs[0]
        p.text = "A CoE provides the structure to scale AI responsibly while maintaining regulatory compliance"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.PRIMARY
        p.alignment = PP_ALIGN.CENTER

    def add_value_prop_slide(self):
        """Slide 3: Value Proposition."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_header(slide, "Strategic Value Proposition")

        # Four value boxes
        values = [
            ("Accelerate Innovation", "50% faster", "time-to-production", self.PRIMARY),
            ("Reduce Risk", "Zero critical", "regulatory findings", self.SUCCESS),
            ("Optimize Investment", "30% reduction", "in duplicated efforts", self.SECONDARY),
            ("Build Capability", "100% coverage", "AI literacy program", self.ACCENT)
        ]

        positions = [(0.5, 1.5), (6.7, 1.5), (0.5, 4.2), (6.7, 4.2)]

        for i, (title, metric, desc, color) in enumerate(values):
            x, y = positions[i]

            # Card
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(6), Inches(2.4))
            card.fill.solid()
            card.fill.fore_color.rgb = self.LIGHT
            card.line.color.rgb = color
            card.line.width = Pt(3)

            # Title
            t_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.2), Inches(5.4), Inches(0.5))
            tf = t_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = color

            # Metric
            m_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.7), Inches(5.4), Inches(0.8))
            tf = m_box.text_frame
            p = tf.paragraphs[0]
            p.text = metric
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = self.DARK

            # Description
            d_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 1.5), Inches(5.4), Inches(0.5))
            tf = d_box.text_frame
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(100, 100, 100)

    def add_structure_slide(self):
        """Slide 4: CoE Structure at a Glance."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_header(slide, "CoE Structure at a Glance")

        # Core Team box
        core_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(6), Inches(3.2))
        core_box.fill.solid()
        core_box.fill.fore_color.rgb = self.PRIMARY
        core_box.line.fill.background()

        core_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.5))
        tf = core_title.text_frame
        p = tf.paragraphs[0]
        p.text = "Core Team: 5 FTEs"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.WHITE

        roles = [
            "CoE Lead - Strategy & Governance",
            "AI Architect - Platform & Standards",
            "Model Risk Specialist - Compliance",
            "AI Product Manager - Hub & Pipeline",
            "AI Trainer - Capability Building"
        ]

        y = 2.1
        for role in roles:
            r_box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(5.5), Inches(0.4))
            tf = r_box.text_frame
            p = tf.paragraphs[0]
            p.text = "• " + role
            p.font.size = Pt(14)
            p.font.color.rgb = self.WHITE
            y += 0.45

        # Champion Network box
        champ_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4), Inches(6), Inches(3.2))
        champ_box.fill.solid()
        champ_box.fill.fore_color.rgb = self.SECONDARY
        champ_box.line.fill.background()

        champ_title = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(0.5))
        tf = champ_title.text_frame
        p = tf.paragraphs[0]
        p.text = "Federated Champions: 6 BUs"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.WHITE

        champions = [
            "Retail Banking Champion",
            "Corporate Banking Champion",
            "Treasury Champion",
            "Operations Champion",
            "Risk & Compliance Champion",
            "Digital Channels Champion"
        ]

        y = 2.1
        for champ in champions:
            c_box = slide.shapes.add_textbox(Inches(7), Inches(y), Inches(5.5), Inches(0.4))
            tf = c_box.text_frame
            p = tf.paragraphs[0]
            p.text = "• " + champ
            p.font.size = Pt(14)
            p.font.color.rgb = self.WHITE
            y += 0.4

        # Key point
        key_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.8))
        key_box.fill.solid()
        key_box.fill.fore_color.rgb = self.LIGHT
        key_box.line.fill.background()

        key_title = slide.shapes.add_textbox(Inches(0.7), Inches(4.95), Inches(12), Inches(0.4))
        tf = key_title.text_frame
        p = tf.paragraphs[0]
        p.text = "Operating Model: Lean Core, Federated Reach"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.PRIMARY

        key_points = slide.shapes.add_textbox(Inches(0.7), Inches(5.4), Inches(12), Inches(1))
        tf = key_points.text_frame
        p = tf.paragraphs[0]
        p.text = "• Small central team sets standards, manages platform, ensures compliance"
        p.font.size = Pt(14)
        p.font.color.rgb = self.DARK
        p2 = tf.add_paragraph()
        p2.text = "• Champions embedded in BUs drive adoption without creating bottlenecks"
        p2.font.size = Pt(14)
        p2.font.color.rgb = self.DARK
        p3 = tf.add_paragraph()
        p3.text = "• Tiered governance: Tier 3 self-service, Tier 2 standard, Tier 1 strategic"
        p3.font.size = Pt(14)
        p3.font.color.rgb = self.DARK

    def add_ownership_slide(self):
        """Slide 5: What the CoE Owns."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_header(slide, "CoE Ownership & Accountability")

        # Six ownership boxes in 2x3 grid
        areas = [
            ("AI Governance", "Framework, policy, steering committee", self.PRIMARY),
            ("Model Risk", "Inventory, validation, monitoring", self.SECONDARY),
            ("MLOps Platform", "Hub, pipelines, templates", self.ACCENT),
            ("Copilot Governance", "O365 Copilot policies & training", self.SUCCESS),
            ("Specialized AI", "Fraud, Agentic AI, high-risk models", RGBColor(156, 39, 176)),
            ("Capability Building", "Training, certification, CoP", RGBColor(0, 150, 136))
        ]

        positions = [
            (0.5, 1.4), (4.5, 1.4), (8.5, 1.4),
            (0.5, 4.0), (4.5, 4.0), (8.5, 4.0)
        ]

        for i, (title, desc, color) in enumerate(areas):
            x, y = positions[i]

            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(2.3))
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.fill.background()

            t_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.3), Inches(3.4), Inches(0.6))
            tf = t_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.WHITE

            d_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.0), Inches(3.4), Inches(1))
            tf = d_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(14)
            p.font.color.rgb = self.WHITE

    def add_investment_slide(self):
        """Slide 6: Investment & ROI."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_header(slide, "Investment & Expected Returns")

        # Left - Investment
        inv_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(5.8), Inches(3))
        inv_box.fill.solid()
        inv_box.fill.fore_color.rgb = self.LIGHT
        inv_box.line.color.rgb = self.PRIMARY
        inv_box.line.width = Pt(2)

        inv_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.4), Inches(0.5))
        tf = inv_title.text_frame
        p = tf.paragraphs[0]
        p.text = "Investment Required"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.PRIMARY

        investments = [
            ("Core Team (5 FTEs)", "Year 1 & ongoing"),
            ("MLOps Platform", "One-time + maintenance"),
            ("AI/Automation Hub", "Platform development"),
            ("Training Programs", "Content & delivery"),
            ("Copilot Licenses", "Based on adoption")
        ]

        y = 2.1
        for item, timing in investments:
            i_box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(5.4), Inches(0.4))
            tf = i_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.DARK
            y += 0.45

        # Right - Returns
        ret_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4), Inches(6), Inches(3))
        ret_box.fill.solid()
        ret_box.fill.fore_color.rgb = self.LIGHT
        ret_box.line.color.rgb = self.SUCCESS
        ret_box.line.width = Pt(2)

        ret_title = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.6), Inches(0.5))
        tf = ret_title.text_frame
        p = tf.paragraphs[0]
        p.text = "Expected Returns (Year 2)"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.SUCCESS

        returns = [
            "30+ production AI use cases",
            "10,000+ hours saved via automation",
            "40% process automation coverage",
            "Zero critical regulatory findings",
            "50% faster model deployment"
        ]

        y = 2.1
        for ret in returns:
            r_box = slide.shapes.add_textbox(Inches(7), Inches(y), Inches(5.6), Inches(0.4))
            tf = r_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"✓ {ret}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.DARK
            y += 0.45

        # Bottom - Payback
        pay_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2))
        pay_box.fill.solid()
        pay_box.fill.fore_color.rgb = self.PRIMARY
        pay_box.line.fill.background()

        pay_title = slide.shapes.add_textbox(Inches(0.7), Inches(4.9), Inches(12), Inches(0.5))
        tf = pay_title.text_frame
        p = tf.paragraphs[0]
        p.text = "Business Case Summary"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.WHITE

        metrics = [
            ("Efficiency Gains", "Direct cost savings from automation and reduced manual effort"),
            ("Risk Reduction", "Avoided regulatory fines and reputational damage"),
            ("Revenue Enablement", "Faster time-to-market for AI-driven products and services")
        ]

        y = 5.4
        for metric, desc in metrics:
            m_box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(12), Inches(0.35))
            tf = m_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {metric}: {desc}"
            p.font.size = Pt(13)
            p.font.color.rgb = self.WHITE
            y += 0.4

    def add_roadmap_slide(self):
        """Slide 7: High-Level Roadmap."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_header(slide, "Implementation Roadmap")

        # Timeline phases
        phases = [
            ("Q1 2025", "Foundation", ["Charter approval", "Core team onboarding", "Governance framework"], self.PRIMARY),
            ("Q2 2025", "Launch", ["Hub v1.0 live", "Champion network", "First use cases"], self.SECONDARY),
            ("Q3-Q4 2025", "Scale", ["20+ use cases", "Training rollout", "Self-service model"], self.ACCENT),
            ("2026", "Optimize", ["30+ use cases", "MLOps Level 3", "Continuous improvement"], self.SUCCESS)
        ]

        x = 0.5
        for quarter, phase, items, color in phases:
            # Phase box
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(3), Inches(4.8))
            box.fill.solid()
            box.fill.fore_color.rgb = self.LIGHT
            box.line.color.rgb = color
            box.line.width = Pt(3)

            # Quarter header
            q_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.5), Inches(3), Inches(0.6))
            q_box.fill.solid()
            q_box.fill.fore_color.rgb = color
            q_box.line.fill.background()

            q_text = slide.shapes.add_textbox(Inches(x), Inches(1.55), Inches(3), Inches(0.5))
            tf = q_text.text_frame
            p = tf.paragraphs[0]
            p.text = quarter
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.WHITE
            p.alignment = PP_ALIGN.CENTER

            # Phase name
            ph_text = slide.shapes.add_textbox(Inches(x + 0.15), Inches(2.2), Inches(2.7), Inches(0.5))
            tf = ph_text.text_frame
            p = tf.paragraphs[0]
            p.text = phase
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = color

            # Items
            y = 2.7
            for item in items:
                i_text = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y), Inches(2.7), Inches(0.4))
                tf = i_text.text_frame
                p = tf.paragraphs[0]
                p.text = "• " + item
                p.font.size = Pt(13)
                p.font.color.rgb = self.DARK
                y += 0.45

            x += 3.2

        # Arrow
        arrow_text = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4))
        tf = arrow_text.text_frame
        p = tf.paragraphs[0]
        p.text = "Foundation  ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━▶  AI-Driven Organization"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.PRIMARY
        p.alignment = PP_ALIGN.CENTER

    def add_governance_slide(self):
        """Slide 8: Governance & Risk Model."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_header(slide, "Risk-Based Governance Model")

        # Three tier boxes
        tiers = [
            ("Tier 1: High Risk", "Steering Committee Approval",
             ["Customer-facing AI", "Credit/fraud models", "Agentic AI", "Regulatory impact"],
             RGBColor(192, 57, 43)),
            ("Tier 2: Medium Risk", "CoE Lead Approval",
             ["Internal analytics", "Process automation", "Standard ML models"],
             self.ACCENT),
            ("Tier 3: Low Risk", "Champion Approval",
             ["Dashboards", "Simple automation", "Approved templates"],
             self.SUCCESS)
        ]

        x = 0.5
        for title, approval, examples, color in tiers:
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.4), Inches(4), Inches(5.3))
            box.fill.solid()
            box.fill.fore_color.rgb = self.LIGHT
            box.line.color.rgb = color
            box.line.width = Pt(3)

            # Header
            h_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.4), Inches(4), Inches(0.7))
            h_box.fill.solid()
            h_box.fill.fore_color.rgb = color
            h_box.line.fill.background()

            t_text = slide.shapes.add_textbox(Inches(x + 0.1), Inches(1.5), Inches(3.8), Inches(0.5))
            tf = t_text.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.WHITE
            p.alignment = PP_ALIGN.CENTER

            # Approval
            a_text = slide.shapes.add_textbox(Inches(x + 0.15), Inches(2.2), Inches(3.7), Inches(0.5))
            tf = a_text.text_frame
            p = tf.paragraphs[0]
            p.text = approval
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = color

            # Examples
            e_title = slide.shapes.add_textbox(Inches(x + 0.15), Inches(2.7), Inches(3.7), Inches(0.4))
            tf = e_title.text_frame
            p = tf.paragraphs[0]
            p.text = "Examples:"
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(100, 100, 100)

            y = 3.1
            for ex in examples:
                ex_text = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y), Inches(3.7), Inches(0.35))
                tf = ex_text.text_frame
                p = tf.paragraphs[0]
                p.text = "• " + ex
                p.font.size = Pt(13)
                p.font.color.rgb = self.DARK
                y += 0.4

            x += 4.25

        # Bottom note
        note = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4))
        tf = note.text_frame
        p = tf.paragraphs[0]
        p.text = "Proportionate governance: Higher risk = more oversight; Lower risk = faster delivery"
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = self.DARK
        p.alignment = PP_ALIGN.CENTER

    def add_decisions_slide(self):
        """Slide 9: Key Decisions Required."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_header(slide, "Decisions Required from Executive Management")

        decisions = [
            ("1", "Charter Approval", "Approve CoE Charter and mandate", "CDO, CRO, CCO"),
            ("2", "Resource Allocation", "Approve 5 FTE core team budget", "CDO, CFO"),
            ("3", "Champion Nomination", "Nominate BU Champions (6 business units)", "BU Heads"),
            ("4", "Platform Investment", "Approve Hub and MLOps platform investment", "CDO, CTO"),
            ("5", "Steering Committee", "Confirm Steering Committee composition", "CDO")
        ]

        y = 1.5
        for num, title, desc, owner in decisions:
            # Number circle
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y), Inches(0.6), Inches(0.6))
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.ACCENT
            circle.line.fill.background()

            num_text = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.1), Inches(0.6), Inches(0.5))
            tf = num_text.text_frame
            p = tf.paragraphs[0]
            p.text = num
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.WHITE
            p.alignment = PP_ALIGN.CENTER

            # Decision box
            d_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(y), Inches(11.5), Inches(0.9))
            d_box.fill.solid()
            d_box.fill.fore_color.rgb = self.LIGHT
            d_box.line.fill.background()

            # Title
            t_text = slide.shapes.add_textbox(Inches(1.5), Inches(y + 0.1), Inches(3), Inches(0.4))
            tf = t_text.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.PRIMARY

            # Description
            desc_text = slide.shapes.add_textbox(Inches(4.5), Inches(y + 0.1), Inches(5), Inches(0.4))
            tf = desc_text.text_frame
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(14)
            p.font.color.rgb = self.DARK

            # Owner
            o_text = slide.shapes.add_textbox(Inches(9.5), Inches(y + 0.1), Inches(3), Inches(0.4))
            tf = o_text.text_frame
            p = tf.paragraphs[0]
            p.text = owner
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.SECONDARY

            y += 1.05

        # Timeline note
        timeline = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8))
        timeline.fill.solid()
        timeline.fill.fore_color.rgb = self.PRIMARY
        timeline.line.fill.background()

        t_text = slide.shapes.add_textbox(Inches(0.7), Inches(6.45), Inches(12), Inches(0.5))
        tf = t_text.text_frame
        p = tf.paragraphs[0]
        p.text = "Target: All decisions approved by end of Q1 2025 to enable Q2 launch"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.WHITE
        p.alignment = PP_ALIGN.CENTER

    def add_summary_slide(self):
        """Slide 10: Executive Summary."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_header(slide, "Executive Summary")

        # Key points
        points = [
            ("Problem", "Fragmented AI initiatives, inconsistent governance, regulatory risk, skills gap"),
            ("Solution", "AI/ML Center of Excellence with lean core team and federated champion network"),
            ("Investment", "5 FTE core team + platform investment (detailed budget in appendix)"),
            ("Returns", "30+ AI use cases, 50% faster delivery, zero critical findings, 40% automation"),
            ("Ask", "Charter approval, resource allocation, champion nominations by Q1 2025")
        ]

        y = 1.5
        for label, content in points:
            # Label box
            l_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(2), Inches(0.8))
            l_box.fill.solid()
            l_box.fill.fore_color.rgb = self.PRIMARY
            l_box.line.fill.background()

            l_text = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.2), Inches(2), Inches(0.5))
            tf = l_text.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.WHITE
            p.alignment = PP_ALIGN.CENTER

            # Content
            c_text = slide.shapes.add_textbox(Inches(2.7), Inches(y + 0.15), Inches(10), Inches(0.7))
            tf = c_text.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = content
            p.font.size = Pt(16)
            p.font.color.rgb = self.DARK

            y += 1.0

        # Bottom CTA
        cta = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(6.0), Inches(8.3), Inches(1.1))
        cta.fill.solid()
        cta.fill.fore_color.rgb = self.ACCENT
        cta.line.fill.background()

        cta_text = slide.shapes.add_textbox(Inches(2.7), Inches(6.15), Inches(8), Inches(0.8))
        tf = cta_text.text_frame
        p = tf.paragraphs[0]
        p.text = "Recommendation: Approve CoE Charter to enable"
        p.font.size = Pt(18)
        p.font.color.rgb = self.WHITE
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = "responsible AI adoption at scale"
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = self.WHITE
        p2.alignment = PP_ALIGN.CENTER

    def add_closing_slide(self):
        """Slide 11: Q&A / Thank You."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.PRIMARY
        bg.line.fill.background()

        # Thank you
        title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "Thank You"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.WHITE
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.6))
        tf = sub.text_frame
        p = tf.paragraphs[0]
        p.text = "Questions & Discussion"
        p.font.size = Pt(28)
        p.font.color.rgb = self.ACCENT
        p.alignment = PP_ALIGN.CENTER

        # Contact
        contact = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5))
        tf = contact.text_frame
        p = tf.paragraphs[0]
        p.text = "AI Engineering Team | Innovation & Digitization"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(180, 180, 180)
        p.alignment = PP_ALIGN.CENTER

    def generate(self):
        """Generate the executive summary presentation."""
        print("Generating Executive Summary presentation...")

        self.add_title_slide()              # 1
        self.add_why_slide()                # 2
        self.add_value_prop_slide()         # 3
        self.add_structure_slide()          # 4
        self.add_ownership_slide()          # 5
        self.add_investment_slide()         # 6
        self.add_roadmap_slide()            # 7
        self.add_governance_slide()         # 8
        self.add_decisions_slide()          # 9
        self.add_summary_slide()            # 10
        self.add_closing_slide()            # 11

        self.prs.save(self.output_path)
        print(f"✓ Presentation saved: {self.output_path}")
        print(f"  File size: {self.output_path.stat().st_size / 1024:.1f} KB")
        print(f"  Total slides: {len(self.prs.slides)}")


def main():
    output_path = Path('/Users/govind/AIEngineering/governance-docs-word/12-AI-ML-CoE-Executive-Summary.pptx')
    generator = ExecPresentationGenerator(output_path)
    generator.generate()
    print("\nExecutive presentation complete!")


if __name__ == '__main__':
    main()
